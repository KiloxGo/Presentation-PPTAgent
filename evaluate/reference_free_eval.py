import os
import json
import pptx
import base64
import argparse
from PIL import Image
from concurrent.futures import ThreadPoolExecutor, as_completed

from openai import OpenAI

client = OpenAI(
    api_key="sk-iS1MfXEJl0of8xSx8eX5ED6v11KkQbU1Z4CV4jfdDMQv57a8",
    base_url="https://api.chatanywhere.tech/v1"
)

# 定义存储评估标准的文件路径
METRIC_FILE = 'metric_dict.json'

# 默认评估标准（保持不变）
DEFAULT_METRIC_DICT = {
    "text": {"criteria": "For the title slide, the text should be concise, precise, and accurately reflect the academic topic, with minimal or no detailed content required. For non-title slides, the main content must be clear, rigorous, and logically structured, using proper headings and well-distinguished levels. Avoid colloquialism, ensure high relevance to the academic subject, and support key points with evidence or clear referencing where necessary. Use consistent and readable font sizes, styles, and colors; text should include appropriate academic terminology and, if relevant, proper citations. Slides may include a Q&A section (it can only be located at the end of the presentation and is optional), provided its content is directly relevant to the academic conference topic and contributes meaningfully to the discussion. Distinguish between colloquial abbreviations and academic abbreviations, using the latter appropriately and avoiding informal or slang-based abbreviations. Text must adapt to background characteristics (e.g., avoiding overlap with complex background elements) and not cross inappropriate visual boundaries that cause disharmony. Specifically, text MUST NOT span across boundaries between different background color blocks or regions (e.g., half the text on a blue background and half on a white background), as this creates visual confusion and disrupts readability.如果有内容被截断直接给60分.", "scale": 100},
    "image": {"criteria": "For the title slide, images, charts, or icons are optional and may be absent without penalty. For non-title slides, images, charts, and icons must be high quality (clear, no distortion, watermarks, or unnecessary borders), appropriately sized, and scientifically relevant to the topic. Images and icons should enhance or explain the academic content (e.g., research data, diagrams, statistical figures). Academic images or charts should have a concise title or caption and clearly indicate the data source if applicable. Non-academic illustrations or icons are acceptable (as long as they align with the text and do not need to significantly enhance or explain academic content) and do not require a title or caption. The inclusion of icons does not affect the scoring. The layout must align with the page style and maintain an academic tone. Get 80 scores for slides without images or icons, except for the title slide where no penalty applies for their absence. No penalty for images or icons containing academic conference logos or trademarks, even if they are more related to branding and affiliations.", "scale": 100},
    "layout": {"criteria": "All elements (text, images, formulas, icons, tables) must be fully visible and entirely within slide boundaries. Text must not cross inappropriate background boundaries or overlap with background elements in a confusing way, ensuring harmony with the background layout.特别检查文字是否横跨背景中不同颜色/区域的交界处（例如：文字一半在蓝色区块、一半在白色区块），这种情况会严重破坏布局和谐性。 Images must not be truncated, stretched, or cropped in a way that obscures key content (e.g., missing portions of graphs, icons, text, or visual details). Elements should be well-aligned, evenly spaced, with clear content hierarchy for easy comprehension. Minimizing unnecessary white space is critical—utilize the entire slide area effectively while strictly avoiding overcrowding. Visual focus should highlight key points, supporting natural reading flow. Follow academic or institutional template standards if available. Severe penalties apply for excessive/unnecessary blank spaces, as well as any image truncation (even within boundaries) or overflow beyond slide edges.如果有内容被截断直接给60分.", "scale": 100},
    "color": {"criteria": "Use high-contrast color specifically between text and background to ensure legibility. Text color must adapt to background characteristics (e.g., dark text on light backgrounds, light text on dark backgrounds) to avoid visual clash and ensure harmony. Avoid overly bright or saturated colors; color schemes should be professional and suitable for academic conferences. Charts and diagrams should use distinguishable colors for different data sets and be considerate of color accessibility. Highlight important information subtly without disrupting the academic style.", "scale": 100},
}

# 读取保存的评估标准（保持不变）
if os.path.exists(METRIC_FILE):
    try:
        with open(METRIC_FILE, 'r', encoding='utf-8') as f:
            METRIC_DICT = json.load(f)
        for key in METRIC_DICT:
            if key not in DEFAULT_METRIC_DICT or "criteria" not in METRIC_DICT[key] or "scale" not in METRIC_DICT[key]:
                print(f"警告: {METRIC_FILE} 中的 {key} 格式不正确，使用默认值。")
                METRIC_DICT[key] = DEFAULT_METRIC_DICT[key]
    except (json.JSONDecodeError, UnicodeDecodeError) as e:
        print(f"读取 {METRIC_FILE} 文件时发生错误 ({e})，将使用默认评估标准。")
        METRIC_DICT = DEFAULT_METRIC_DICT.copy()
else:
    METRIC_DICT = DEFAULT_METRIC_DICT.copy()

# 核心修改：优化提示词，要求结构化输出
INSTRUCTION = """Please evaluate the slide based on the following criteria:
{}
Provide an integer score between 0-{} (higher scores indicate better adherence to the criteria). 

Your response must strictly follow this format:
1. First line: Only the integer score (e.g., "80")
2. Second line: "[问题] 具体描述问题位置和内容（例如：第3页"核心优势"部分有3行冗余描述："这是我们产品的主要特点..."，第3页图表右侧被截断，缺失数据标签；图片左侧边缘人物面部不完整，第3页"研究方法"标题横跨蓝色与白色背景交界处；段落文字一半在灰色区块外）"
3. Third line: "[修改方向] 具体改进操作（结合Slidev特性，例如：删除冗余修饰，使用Slidev列表组件<List><Item>...</Item></List>，调整图片尺寸至完整显示，使用<Image fit="contain" />确保内容不截断）"
4. Fourth line: "[参考示例] 具体可复用的Slidev代码或格式（例如：<List><Item>高效：比传统工具快30%</Item></List>）"

Notes:
- 【问题】对于图片截断，必须明确位置（如"图表右侧边缘"、"图片底部"）和截断内容（如"缺失坐标轴标签"、"图标一半被切掉"），对于文字横跨背景色块的情况，必须明确位置（如"标题行"、"第三段文字"）和横跨的具体背景区域（如"红蓝背景交界处"、"灰色区块边缘"）
- 【修改方向】需关联Slidev特性（如列表组件<List>、网格布局<Grid>、居中类mx-auto等），针对图像，需关联Slidev的图片显示特性（如fit属性控制缩放方式、调整尺寸避免截断），需关联Slidev的布局特性（如<Box>组件限定背景、调整margin/padding避免跨界、使用背景色块包裹文字）
- 【参考示例】需提供可直接使用的Slidev语法（如组件标签、CSS类），需提供防止截断的Slidev图片组件用法（如指定fit="contain"或调整width/height），需提供避免文字跨界的Slidev语法（如带背景色的容器组件、定位样式）
- If there are no issues, use "[问题] 无明显问题"、"[修改方向] 无需修改"、"[参考示例] 保持当前格式"
- Strictly limit to the criteria, do not evaluate other aspects.
- Use the full scoring range (1-100), not just 60-100.
"""

def encode_image(image_path):
    with open(image_path, 'rb') as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def evaluate_slide(prompt: str, image_url: str, model_name: str, max_tokens: int) -> str:
    messages = [{
        "role": "user",
        "content": [
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": image_url}},
        ],
    }]
    try:
        response = client.chat.completions.create(
            model=model_name,
            messages=messages,
            max_tokens=max_tokens,
            temperature=0.0,
        )
        response_text = response.choices[0].message.content.strip()
        lines = response_text.split('\n')
        # 解析结构化输出
        score = float(lines[0].strip()) if len(lines)>=1 else 0.0
        problem = lines[1].strip() if len(lines)>=2 else "[问题] 解析错误"
        direction = lines[2].strip() if len(lines)>=3 else "[修改方向] 解析错误"
        example = lines[3].strip() if len(lines)>=4 else "[参考示例] 解析错误"
        return {
            "score": score,
            "problem": problem,
            "modification_direction": direction,
            "reference_example": example
        }
    except Exception as e:
        print(f"评估幻灯片时出错: {e}")
        return {
            "score": 0.0,
            "problem": "[问题] 评估出错",
            "modification_direction": "[修改方向] 无法生成建议",
            "reference_example": "[参考示例] 无"
        }

def parse_text(path, page: int = 0):
    slide = pptx.Presentation(path).slides[page]
    text_list = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_list.append(shape.text_frame.text)
    return [t for t in text_list if t.strip()]

def main():
    if args.update_metric:
        name, criteria, scale = args.update_metric
        try:
            scale = int(scale)
            if 0 <= scale <= 100:
                METRIC_DICT[name] = {"criteria": criteria, "scale": scale}
                try:
                    with open(METRIC_FILE, 'w', encoding='utf-8') as f:
                        json.dump(METRIC_DICT, f, indent=4, ensure_ascii=False)
                    print(f"评估标准已更新并保存到 {METRIC_FILE} 文件中。")
                except Exception as e:
                    print(f"保存评估标准到 {METRIC_FILE} 文件时出错: {e}")
            else:
                print(f"错误: scale 必须在 0-100 之间，当前值为 {scale}，使用默认值。")
                METRIC_DICT[name] = DEFAULT_METRIC_DICT.get(name, {"criteria": "", "scale": 100})
        except ValueError:
            print(f"错误: scale 必须是整数，当前值为 {scale}，使用默认值。")
            METRIC_DICT[name] = DEFAULT_METRIC_DICT.get(name, {"criteria": "", "scale": 100})

    if args.image_path:
        if args.image_path.endswith(".png"):
            jpg_image_path = args.image_path.replace('.png', '.jpg')
            if os.path.exists(jpg_image_path):
                print("JPG image already exists! Skip conversion.")
            else:
                png_image = Image.open(args.image_path)
                rgb_image = png_image.convert("RGB")
                rgb_image.save(jpg_image_path, format="JPEG", quality=95)
                print("Conversion successful! Saved as ", jpg_image_path)
        else:
            jpg_image_path = args.image_path

        image_base64 = encode_image(jpg_image_path)
        image_url = f"data:image/jpeg;base64,{image_base64}"

        results_dict = {}
        evaluation_tasks = []
        for metric, metric_dict in METRIC_DICT.items():
            if metric not in args.metric_to_use:
                continue
            prompt = INSTRUCTION.format(metric_dict["criteria"], metric_dict["scale"])
            if metric == "text":
                pptx_path = jpg_image_path.replace(".jpg", ".pptx")
                if os.path.exists(pptx_path):
                    slide_texts = parse_text(pptx_path)
                    prompt += '\n\nTexts in the slide are: \n' + '\n'.join(slide_texts)
                else:
                    print(f"警告: 未找到对应的 PPTX 文件 {pptx_path}，跳过文本解析。")
            evaluation_tasks.append((metric, prompt))

        with ThreadPoolExecutor(max_workers=len(evaluation_tasks)) as executor:
            future_to_metric = {
                executor.submit(evaluate_slide, prompt, image_url, args.model_name, args.max_tokens): metric
                for metric, prompt in evaluation_tasks
            }
            for future in as_completed(future_to_metric):
                metric = future_to_metric[future]
                try:
                    result = future.result()
                    results_dict[metric] = result
                except Exception as e:
                    print(f"评估维度 {metric} 时出错: {e}")
                    results_dict[metric] = {
                        "score": 0.0,
                        "problem": "[问题] 并行评估出错",   
                        "modification_direction": "[修改方向] 无",
                        "reference_example": "[参考示例] 无"
                    }

        for metric, result in results_dict.items():
            print(f"{metric} 评分: {result['score']}")
            print(f"问题: {result['problem']}")
            print(f"修改方向: {result['modification_direction']}")
            print(f"参考示例: {result['reference_example']}\n")
        print("总分: ", sum([v["score"] for v in results_dict.values()]))

        if args.response_path is None:
            response_path = jpg_image_path.replace(".jpg", "_eval.json")
            with open(response_path, 'w', encoding='utf-8') as f:
                json.dump(results_dict, f, indent=4, ensure_ascii=False)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX reference-free evaluation.")
    parser.add_argument("--image_path", type=str, default=None,
                        help="Path to the image of model-generated PPTX file.")
    parser.add_argument("--response_path", type=str, default=None,
                        help="Path to save the model-evaluation justifications.")
    parser.add_argument("--metric_to_use", type=str, nargs='+',
                        default=["text", "image", "layout", "color"],
                        help="Metrics to evaluate.")
    parser.add_argument("--model_name", type=str, default="gpt-4o", help="Model name to use.")
    parser.add_argument("--max_tokens", type=int, default=300,  # 增加token上限以容纳结构化内容
                        help="Max tokens to generate.")
    parser.add_argument("--update_metric", type=str, nargs=3,
                        default=None, help="Update METRIC_DICT with new metric. Format: name criteria scale")

    args = parser.parse_args()

    media_dir = os.path.join(os.path.dirname(args.image_path), "media") if args.image_path else None
    if media_dir and not os.path.exists(media_dir):
        print("[Warning] No image is involved in this example. May affect the interpretation of the image criteria.")

    main()