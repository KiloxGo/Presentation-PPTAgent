from playwright.sync_api import sync_playwright
import time
import os
import argparse
import subprocess
import socket
from multiprocessing import Pool
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def get_free_port():
    """获取一个空闲的端口（解决端口冲突问题）"""
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(('localhost', 0))  # 绑定到0端口，系统会自动分配空闲端口
        return s.getsockname()[1]  # 返回分配的端口号

def process_single_file(args):
    """处理单个md文件（适配多进程，接收元组参数：(md_file, output_root)）"""
    md_file, output_root = args  # 多进程传递参数需要打包成元组
    file_name = os.path.splitext(os.path.basename(md_file))[0]
    output_dir = output_root
    os.makedirs(output_dir, exist_ok=True)
    
    img_path = os.path.join(output_dir, f"{file_name}.jpg")
    ppt_path = os.path.join(output_dir, f"{file_name}.pptx")
    
    # 动态获取空闲端口
    port = get_free_port()
    slidev_process = None  # 用于管理当前Slidev进程

    # 启动浏览器并截图
    with sync_playwright() as p:
        try:
            print(f"\n正在处理: {md_file}（端口：{port}）")
            # 启动Slidev服务（指定当前文件和动态端口）
            slidev_cmd = f"slidev {md_file} --port {port}"
            slidev_process = subprocess.Popen(slidev_cmd, shell=True)
            time.sleep(15)  # 等待服务启动（可根据实际情况调整）
            
            # 启动浏览器访问当前端口的服务
            browser = p.chromium.launch(headless=True)
            page = browser.new_page()
            page.goto(f"http://localhost:{port}")
            
            # 等待页面加载完成
            time.sleep(6)
            
            # 截图第一页
            print(f"截取{md_file}页面...")
            page.screenshot(path=img_path, full_page=True)
            browser.close()
            print(f"截图已保存至: {os.path.abspath(img_path)}")
            
        except Exception as e:
            print(f"处理{md_file}时出错: {str(e)}")
            return False
        finally:
            # 精确关闭当前Slidev进程（只杀当前实例，不影响其他进程）
            if slidev_process:
                print(f"关闭{md_file}的Slidev服务（端口：{port}）...")
                slidev_process.terminate()  # 终止当前进程
                slidev_process.wait(timeout=5)  # 等待进程退出
            time.sleep(5)  # 确保端口释放
    
    # 生成PPT（逻辑不变）
    try:
        print(f"生成{md_file}的PPT...")
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9比例
        prs.slide_height = Inches(7.5)
        
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        with Image.open(img_path) as img:
            img_width_px, img_height_px = img.size
            img_ratio = img_width_px / img_height_px
        
        slide_ratio = slide_width / slide_height
        if img_ratio > slide_ratio:
            display_width = slide_width
            display_height = display_width / img_ratio
        else:
            display_height = slide_height
            display_width = display_height * img_ratio
        
        left = (slide_width - display_width) / 2
        top = (slide_height - display_height) / 2
        
        slide.shapes.add_picture(
            img_path,
            left=Inches(left),
            top=Inches(top),
            width=Inches(display_width),
            height=Inches(display_height)
        )
        
        prs.save(ppt_path)
        print(f"PPT已保存至: {os.path.abspath(ppt_path)}")
        return True
        
    except Exception as e:
        print(f"生成{md_file}的PPT出错: {str(e)}")
        return False

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="并行处理Slidev的md文件并导出PPT")
    parser.add_argument(
        "files", 
        nargs="+",
        help="需要处理的Markdown文件路径（可指定多个）"
    )
    parser.add_argument(
        "-o", "--output", 
        type=str, 
        default="custom_exports", 
        help="指定输出根目录（默认：custom_exports）"
    )
    parser.add_argument(
        "-p", "--processes", 
        type=int, 
        default=4, 
        help="并行进程数（默认：4，根据CPU核心数调整）"
    )
    args = parser.parse_args()
    
    # 检查并安装依赖
    try:
        import playwright
        import pptx
        from PIL import Image
    except ImportError as e:
        print("正在安装所需依赖...")
        os.system("pip install playwright python-pptx pillow")
        os.system("playwright install")
    
    # 验证文件有效性
    valid_files = []
    for file_path in args.files:
        if os.path.isfile(file_path) and file_path.lower().endswith(".md"):
            valid_files.append(os.path.abspath(file_path))
        else:
            print(f"跳过无效文件: {file_path}")
    
    if not valid_files:
        print("没有找到有效的Markdown文件，程序退出")
    else:
        print(f"共找到{len(valid_files)}个有效文件，开始并行处理...")
        # 构造多进程参数（每个文件对应一个(output_root, md_file)元组）
        process_args = [(file, args.output) for file in valid_files]
        # 使用进程池并行处理
        with Pool(processes=args.processes) as pool:
            pool.map(process_single_file, process_args)
        print("\n所有文件处理完毕")