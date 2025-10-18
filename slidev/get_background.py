import os
from pptx import Presentation
from pptx.util import Inches
import subprocess

def clear_slide_content(pptx_path, temp_pptx_path):
    size_limit_inch = 50  # 只删除宽高都小于这个英寸值的形状
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in list(slide.shapes):
            # 删除文字框内容
            if shape.has_text_frame:
                shape.text_frame.clear()
            # 只删除宽高都小于阈值的图片、图标、表格、嵌入对象、SmartArt、图表
            elif shape.shape_type in [13, 75, 5, 30, 21, 3]:
                if shape.width < Inches(size_limit_inch) and shape.height < Inches(size_limit_inch):
                    try:
                        shape._element.getparent().remove(shape._element)
                    except Exception as e:
                        print(f"无法删除shape，类型{shape.shape_type}，错误: {e}")
            else:
                # 其他类型直接强制删除
                try:
                    shape._element.getparent().remove(shape._element)
                except Exception as e:
                    print(f"无法删除其他类型shape，类型{shape.shape_type}，错误: {e}")
    prs.save(temp_pptx_path)

def export_slide_with_libreoffice(temp_pptx_path, output_folder, idx, slide_idx):
    prs = Presentation(temp_pptx_path)
    for i in reversed(range(len(prs.slides))):
        if i != slide_idx:
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])
    single_slide_pptx = os.path.join(output_folder, f"single_slide_{idx}.pptx")
    prs.save(single_slide_pptx)

    cmd = [
        "soffice",
        "--headless",
        "--convert-to", "jpg",
        "--outdir", output_folder,
        single_slide_pptx
    ]
    print(f"正在导出第{idx}页图片...")
    result = subprocess.run(cmd, capture_output=True, text=True)
    basename = os.path.splitext(os.path.basename(single_slide_pptx))[0]
    files = [f for f in os.listdir(output_folder) if f.startswith(basename) and f.endswith(".jpg")]
    if files:
        src = os.path.join(output_folder, files[0])
        dst = os.path.join(output_folder, f"image{idx}.jpg")
        os.rename(src, dst)
        print(f"已保存：{dst}")
    else:
        print("未找到导出的图片文件")
    os.remove(single_slide_pptx)

def main():
    input_pptx = r"E:\prp\background\source.pptx"  # 输入PPTX文件路径
    output_folder = r"E:\prp\background\background_images3"  # 输出图片文件夹路径
    temp_pptx = os.path.join(output_folder, "temp_bg_only.pptx")

    clear_slide_content(input_pptx, temp_pptx)

    prs = Presentation(temp_pptx)
    for idx in range(len(prs.slides)):
        export_slide_with_libreoffice(temp_pptx, output_folder, idx + 1, idx)

    if os.path.exists(temp_pptx):
        os.remove(temp_pptx)
    print("全部导出完成！")

if __name__ =="__main__":
    main()