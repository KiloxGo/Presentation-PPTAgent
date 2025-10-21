import os
import json
import pptx
import base64
import argparse
from PIL import Image

from openai import OpenAI
from config import get_openai_config, get_evaluation_criteria, get_evaluation_metrics, get_metric_file

openai_config = get_openai_config()
client = OpenAI(
    api_key=openai_config['key'],
    base_url=openai_config['baseUrl']
)

# 定义存储评估标准的文件路径
METRIC_FILE = get_metric_file()

# 读取保存的评估标准
if os.path.exists(METRIC_FILE):
    try:
        with open(METRIC_FILE, 'r') as f:
            METRIC_DICT = json.load(f)
    except json.JSONDecodeError:
        print(f"读取 {METRIC_FILE} 文件时发生 JSON 解析错误，将使用默认评估标准。")
        METRIC_DICT = get_evaluation_criteria()
else:
    METRIC_DICT = get_evaluation_criteria()

INSTRUCTION = """Please evaluate the slide based on the following criteria:
{}

Give an integer score between 0 - {}, higher scores means the criteria is better met.
First, respond with a score; Then, provide your justification for the score in natural language sentences. Your response should look like this: '4. The slide has concise texts summarized in bullet points.'
Only evaluate the slide based on the specified criteria, and no other aspects. Give scores across the full spectrum (1-5) instead of only good ones (3-5).
"""

def encode_image(image_path):
    with open(image_path, 'rb') as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def evaluate_slide(prompt: str, image_url: str) -> str:
    messages = [{
        "role": "user",
        "content": [
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": image_url}},
        ],
    }]
    response = client.chat.completions.create(
            model=args.model_name,
            messages=messages,
            max_tokens=args.max_tokens,
            temperature=0.0,
    )
    response = response.choices[0].message.content.strip()
    try:
        score, justification = response.split(".", 1)
        score = float(score.strip())
    except:
        score, justification = 0.0, response
    return {
        "score": score,
        "justification": justification.strip(),
    }

def parse_text(path, page: int = 0):
    slide = pptx.Presentation(path).slides[page]
    text_list = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_list.append(shape.text_frame.text)
    text_list = [t for t in text_list if t.strip()]
    return text_list

def main():
    # 更新 METRIC_DICT 的值
    if args.update_metric:
        name, criteria, scale = args.update_metric
        METRIC_DICT[name] = {"criteria": criteria, "scale": int(scale)}
        # 保存更新后的评估标准
        try:
            with open(METRIC_FILE, 'w') as f:
                json.dump(METRIC_DICT, f, indent=4)
            print(f"评估标准已更新并保存到 {METRIC_FILE} 文件中。")
        except Exception as e:
            print(f"保存评估标准到 {METRIC_FILE} 文件时出错: {e}")

    if args.image_path:
        if args.image_path.endswith(".png"):
            jpg_image_path = args.image_path.replace('.png', '.jpg')
            if os.path.exists(jpg_image_path):
                print("JPG image already exists! Skip conversion.")
            else:
                png_image = Image.open(args.image_path)
                rgb_image = png_image.convert("RGB")
                rgb_image.save(jpg_image_path, format="JPEG", quality=95)
                print("Conversion successful! Saved as ", jpg_image_path)
        else:
            jpg_image_path = args.image_path

        image_base64 = encode_image(jpg_image_path)
        image_url = f"data:image/jpeg;base64,{image_base64}"

        results_dict = {}
        for metric, metric_dict in METRIC_DICT.items():
            if metric not in args.metric_to_use: continue
            prompt = INSTRUCTION.format(metric_dict["criteria"], metric_dict["scale"])
            if metric == "text":
                pptx_path = jpg_image_path.replace(".jpg", ".pptx")
                slide_texts = parse_text(pptx_path)
                prompt += '\n\nTexts in the slide are: \n' + '\n'.join(slide_texts)
            results_dict[metric] = evaluate_slide(prompt, image_url)

        for metric, result in results_dict.items():
            print(f"{metric}: {result}\n")
        print("Total: ", sum([v["score"] for v in results_dict.values()]))

        if args.response_path is None:
            response_path = jpg_image_path.replace(".jpg", "_eval.json")
            json.dump(results_dict, open(response_path, "w"))

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX reference-free evaluation.")
    parser.add_argument("--image_path", type=str, default=None,
                        help="Path to the image of model-generated PPTX file.")
    parser.add_argument("--response_path", type=str, default=None,
                        help="Path to save the model-evaluation justifications.")
    parser.add_argument("--metric_to_use", type=str, nargs='+',
                        default=get_evaluation_metrics(),
                        help="Metrics to evaluate.")
    parser.add_argument("--model_name", type=str, default=openai_config['model'], help="Model name to use.")
    parser.add_argument("--max_tokens", type=int, default=openai_config['maxTokens'], help="Max tokens to generate.")
    # 新增命令行参数用于更新评估标准
    parser.add_argument("--update_metric", type=str, nargs=3,
                        default=None, help="Update METRIC_DICT with new metric. Format: name criteria scale")

    args = parser.parse_args()

    media_dir = os.path.join(os.path.dirname(args.image_path), "media") if args.image_path else None
    if media_dir and not os.path.exists(media_dir):
        print("[Warning] No image is involved in this example. May affect the interpretation of the image criteria.")

    main()